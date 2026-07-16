from __future__ import annotations

import argparse
from pathlib import Path
from typing import Any

import cv2
import numpy as np
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


WIDTH, HEIGHT = 1280, 720
NAVY = "17324D"
TEAL = "008C95"
RED = "C4473A"
GREEN = "2E7D5B"
GOLD = "D39B2A"
INK = "1E2933"
MUTED = "60707D"
PAPER = "F4F7F8"
WHITE = "FFFFFF"


SLIDES = [
    ("LiteTest-MAS", "赛题对齐强化与补充实验", ["多智能体测试生成", "M9 与 M9.1 公开结果", "固定任务 · 固定模型 · 可复核边界"]),
    ("研究问题", "质量、通信与状态如何分开计量？", ["Protocol 是否降低 Agent 通信开销", "StateVector 是否替代重复文本状态", "门控 Memory 是否减少无关注入与负迁移"]),
    ("多 Agent 协作链", "逻辑角色不等于独立 LLM 调用", ["Planner → Retriever → TestGen → Executor → Summarizer", "角色事件记录 action、input/output reference 与状态", "official tests 始终位于 Agent 上下文之外"]),
    ("CompactProtocol V2", "静态信息注册一次，后续只传引用", ["sequence 级一次握手与版本协商", "capability_id、task_ref、reference_id", "确定性序列化与本地可审计 registry"]),
    ("StateVector V2", "固定 bytes 替代同义状态文本", ["网络字节序与 schema 校验", "33 bytes 对 130 bytes 等价文本", "压缩率 0.7462；不等于质量必然提升"]),
    ("GatedSharedMemory V2", "检索、门控、注入与复用分开记录", ["query → candidate → hit → accept/reject/abstain", "按 dataset、task group、seed、experiment 隔离", "有效复用不是因果性能提升"]),
    ("M9：质量提高伴随成本增加", "Protocol V1 提升本次成功率，但不更省模型开销", ["G1/G2/G3/G4 成功率：0.60 / 0.75 / 0.75 / 0.70", "G2-G1：+0.15，95% CI [+0.0333, +0.2671]", "G4 相对 G3：-0.05，保留潜在负迁移"]),
    ("M9.1：V2 补充实验", "S3 最高，S4 负结果仍然存在", ["S1/S2/S3/S4 成功率：0.55 / 0.55 / 0.60 / 0.55", "StateVector V2 记录到真实压缩", "通信 Token 部分 unavailable，不声称 Protocol V2 已证实更省"]),
    ("安全与可复现", "公开层只保留脱敏聚合与校验信息", ["240/240 final records，Strict Verifier valid", "private tests、candidate code、raw response 不公开", "Windows/openEuler 同 SHA 测试与离线 Dashboard"]),
    ("结论", "工程证据完整，性能结论保持克制", ["Protocol 的质量收益不等于通信节省", "StateVector 证明状态压缩，不证明质量必升", "Memory 有真实复用，负迁移风险仍需保留"]),
]


def _rgb(value: str) -> RGBColor:
    """将六位十六进制颜色转换为 PPT RGBColor。"""
    return RGBColor.from_string(value)


def _font(size: int, bold: bool = False) -> ImageFont.FreeTypeFont:
    """加载 Windows 中文字体，用于生成视频帧。"""
    candidates = (
        Path("C:/Windows/Fonts/msyh.ttc"),
        Path("C:/Windows/Fonts/simhei.ttf"),
        Path("/usr/share/fonts/google-noto-cjk/NotoSansCJK-Regular.ttc"),
    )
    for path in candidates:
        if path.is_file():
            index = 1 if bold and path.suffix.lower() == ".ttc" else 0
            return ImageFont.truetype(str(path), size=size, index=index)
    return ImageFont.load_default()


def _add_text(slide: Any, text: str, x: float, y: float, w: float, h: float, size: int, color: str, bold: bool = False) -> None:
    """向 PPT 页面添加无边框文本框。"""
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    paragraph.font.name = "Microsoft YaHei"
    paragraph.font.size = Pt(size)
    paragraph.font.bold = bold
    paragraph.font.color.rgb = _rgb(color)
    paragraph.alignment = PP_ALIGN.LEFT


def _add_header(slide: Any, title: str, subtitle: str, number: int) -> None:
    """添加统一页眉、标题和页码。"""
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.16), Inches(7.5))
    band.fill.solid()
    band.fill.fore_color.rgb = _rgb(TEAL)
    band.line.fill.background()
    _add_text(slide, title, 0.65, 0.55, 11.8, 0.65, 28, NAVY, True)
    _add_text(slide, subtitle, 0.67, 1.22, 11.6, 0.45, 15, MUTED)
    _add_text(slide, f"{number:02d}", 12.25, 0.55, 0.55, 0.35, 11, MUTED, True)


def _add_bullets(slide: Any, bullets: list[str]) -> None:
    """添加稳定尺寸的三条中文要点。"""
    colors = (TEAL, GOLD, RED)
    for index, bullet in enumerate(bullets):
        y = 2.15 + index * 1.25
        marker = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.72), Inches(y), Inches(0.12), Inches(0.72))
        marker.fill.solid()
        marker.fill.fore_color.rgb = _rgb(colors[index])
        marker.line.fill.background()
        _add_text(slide, bullet, 1.05, y - 0.02, 11.4, 0.8, 21, INK, index == 0)


def _add_result_bars(slide: Any, values: list[float], labels: list[str]) -> None:
    """添加任务成功率原生条形图。"""
    colors = (NAVY, TEAL, GOLD, RED)
    for index, (label, value) in enumerate(zip(labels, values)):
        x = 1.1 + index * 2.9
        height = value * 2.2
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(6.7 - height), Inches(1.55), Inches(height))
        bar.fill.solid()
        bar.fill.fore_color.rgb = _rgb(colors[index])
        bar.line.fill.background()
        _add_text(slide, f"{value:.2f}", x + 0.35, 6.35 - height, 0.9, 0.3, 14, INK, True)
        _add_text(slide, label, x + 0.45, 6.78, 0.8, 0.3, 13, MUTED, True)


def build_presentation(output: Path) -> dict[str, int]:
    """生成匿名、可编辑的中文赛事演示 PPTX。"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    for number, (title, subtitle, bullets) in enumerate(SLIDES, 1):
        slide = prs.slides.add_slide(blank)
        background = slide.background.fill
        background.solid()
        background.fore_color.rgb = _rgb(PAPER)
        _add_header(slide, title, subtitle, number)
        if number in (7, 8):
            values = [0.60, 0.75, 0.75, 0.70] if number == 7 else [0.55, 0.55, 0.60, 0.55]
            labels = ["G1", "G2", "G3", "G4"] if number == 7 else ["S1", "S2", "S3", "S4"]
            _add_text(slide, bullets[1], 0.75, 1.95, 11.8, 0.48, 18, INK, True)
            _add_result_bars(slide, values, labels)
        else:
            _add_bullets(slide, bullets)
        notes = slide.notes_slide.notes_text_frame
        notes.text = f"[时长：约 35 秒]\n核心信息：{subtitle}\n过渡：进入下一部分。"
    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output)
    return {"slide_count": len(prs.slides), "bytes": output.stat().st_size}


def _draw_frame(title: str, subtitle: str, bullets: list[str], number: int) -> np.ndarray:
    """绘制与 PPT 内容一致的 16:9 中文视频帧。"""
    image = Image.new("RGB", (WIDTH, HEIGHT), "#F4F7F8")
    draw = ImageDraw.Draw(image)
    draw.rectangle((0, 0, 18, HEIGHT), fill="#008C95")
    draw.text((65, 55), title, font=_font(44, True), fill="#17324D")
    draw.text((68, 120), subtitle, font=_font(25), fill="#60707D")
    draw.text((1190, 60), f"{number:02d}", font=_font(18, True), fill="#60707D")
    colors = ("#008C95", "#D39B2A", "#C4473A")
    for index, bullet in enumerate(bullets):
        y = 225 + index * 125
        draw.rectangle((70, y, 82, y + 68), fill=colors[index])
        draw.text((110, y + 8), bullet, font=_font(29, index == 0), fill="#1E2933")
    draw.text((70, 670), "LiteTest-MAS · 公开聚合演示", font=_font(16), fill="#60707D")
    return cv2.cvtColor(np.asarray(image), cv2.COLOR_RGB2BGR)


def build_video(output: Path, fps: int = 10, seconds_per_slide: int = 4) -> dict[str, int]:
    """从演示内容生成带淡入过渡的无声 MP4。"""
    output.parent.mkdir(parents=True, exist_ok=True)
    writer = cv2.VideoWriter(str(output), cv2.VideoWriter_fourcc(*"mp4v"), fps, (WIDTH, HEIGHT))
    if not writer.isOpened():
        raise RuntimeError("video_writer_unavailable")
    frames = [_draw_frame(title, subtitle, bullets, number) for number, (title, subtitle, bullets) in enumerate(SLIDES, 1)]
    total = 0
    hold = fps * seconds_per_slide
    transition = fps // 2
    for index, frame in enumerate(frames):
        for _ in range(hold):
            writer.write(frame)
            total += 1
        if index + 1 < len(frames):
            next_frame = frames[index + 1]
            for step in range(1, transition + 1):
                alpha = step / (transition + 1)
                writer.write(cv2.addWeighted(frame, 1 - alpha, next_frame, alpha, 0))
                total += 1
    writer.release()
    return {"frame_count": total, "width": WIDTH, "height": HEIGHT, "bytes": output.stat().st_size}


def main() -> int:
    """执行赛事 PPTX 与演示视频构建。"""
    parser = argparse.ArgumentParser(description="生成 LiteTest-MAS 中文赛事演示材料")
    parser.add_argument("--pptx", default="LiteTest-MAS赛事演示.pptx")
    parser.add_argument("--video", default="LiteTest-MAS演示视频.mp4")
    args = parser.parse_args()
    pptx = build_presentation(Path(args.pptx))
    video = build_video(Path(args.video))
    print({"pptx": pptx, "video": video})
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
